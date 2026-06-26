"""Build The Sims 5 Real Estate hackathon PowerPoint deck."""

from __future__ import annotations

import shutil
import subprocess
import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "outputs" / "presentation"
ASSETS = OUT_DIR / "assets"
CURSOR_ASSETS = Path(
    r"C:\Users\kufas\.cursor\projects\c-Users-kufas-OneDrive-Documents-GitHub-The-Sims-5-Real-Estate\assets"
)

IMAGE_NAMES = [
    "ppt-hero-abu-dhabi.png",
    "ppt-tenant-simulator.png",
    "ppt-investment-dashboard.png",
]


def ensure_deps() -> None:
    try:
        import pptx  # noqa: F401
    except ImportError:
        subprocess.check_call(
            [sys.executable, "-m", "pip", "install", "python-pptx", "pillow", "-q"]
        )


def copy_assets() -> None:
    ASSETS.mkdir(parents=True, exist_ok=True)
    for name in IMAGE_NAMES:
        src = CURSOR_ASSETS / name
        dst = ASSETS / name
        if src.exists():
            shutil.copy2(src, dst)
        elif not dst.exists():
            raise FileNotFoundError(f"Missing image: {src}")


def rgb(hex_color: str):
    hex_color = hex_color.lstrip("#")
    return tuple(int(hex_color[i : i + 2], 16) for i in (0, 2, 4))


def add_title_slide(prs, title: str, subtitle: str, image_path: Path | None = None):
    from pptx.util import Inches, Pt

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb("0F172A")

    if image_path and image_path.exists():
        slide.shapes.add_picture(str(image_path), Inches(0), Inches(0), width=Inches(10))
        overlay = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(7.5))
        overlay.fill.solid()
        overlay.fill.fore_color.rgb = rgb("0F172A")
        overlay.fill.transparency = 0.35
        overlay.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.6), Inches(2.0), Inches(8.8), Inches(2.5))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = rgb("F5F0E8")

    p2 = tf.add_paragraph()
    p2.text = subtitle
    p2.font.size = Pt(20)
    p2.font.color.rgb = rgb("C9B896")
    p2.space_before = Pt(12)

    tag = slide.shapes.add_textbox(Inches(0.6), Inches(6.5), Inches(8.8), Inches(0.6))
    tag_tf = tag.text_frame.paragraphs[0]
    tag_tf.text = "Abu Dhabi AI PropTech Challenge · Investment + Decision Intelligence"
    tag_tf.font.size = Pt(14)
    tag_tf.font.color.rgb = rgb("94A3B8")


def add_section_slide(prs, title: str, subtitle: str):
    from pptx.util import Inches, Pt

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb("1E293B")

    accent = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.15), Inches(7.5))
    accent.fill.solid()
    accent.fill.fore_color.rgb = rgb("C9B896")
    accent.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(8.5), Inches(2))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = rgb("F5F0E8")
    p2 = tf.add_paragraph()
    p2.text = subtitle
    p2.font.size = Pt(18)
    p2.font.color.rgb = rgb("94A3B8")
    p2.space_before = Pt(10)


def add_content_slide(
    prs,
    title: str,
    bullets: list[str],
    image_path: Path | None = None,
    image_right: bool = True,
):
    from pptx.util import Inches, Pt

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb("0F172A")

    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(9), Inches(0.8))
    tp = title_box.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(28)
    tp.font.bold = True
    tp.font.color.rgb = rgb("F5F0E8")

    text_left = Inches(0.6)
    text_width = Inches(5.2)
    if image_path and image_path.exists():
        if image_right:
            slide.shapes.add_picture(
                str(image_path), Inches(5.9), Inches(1.2), width=Inches(3.8)
            )
        else:
            slide.shapes.add_picture(
                str(image_path), Inches(0.6), Inches(1.2), width=Inches(3.8)
            )
            text_left = Inches(4.7)
    else:
        text_width = Inches(8.8)

    body = slide.shapes.add_textbox(text_left, Inches(1.2), text_width, Inches(5.8))
    tf = body.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(16 if len(bullets) > 5 else 18)
        p.font.color.rgb = rgb("E2E8F0")
        p.space_after = Pt(10)


def add_two_column_slide(prs, title: str, left_title: str, left_items: list[str], right_title: str, right_items: list[str]):
    from pptx.util import Inches, Pt

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb("0F172A")

    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(9), Inches(0.8))
    tp = title_box.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(28)
    tp.font.bold = True
    tp.font.color.rgb = rgb("F5F0E8")

    for col, col_title, items, x in [
        (0, left_title, left_items, 0.6),
        (1, right_title, right_items, 5.2),
    ]:
        head = slide.shapes.add_textbox(Inches(x), Inches(1.1), Inches(4.2), Inches(0.5))
        hp = head.text_frame.paragraphs[0]
        hp.text = col_title
        hp.font.size = Pt(20)
        hp.font.bold = True
        hp.font.color.rgb = rgb("C9B896")

        body = slide.shapes.add_textbox(Inches(x), Inches(1.7), Inches(4.2), Inches(5.2))
        tf = body.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = item
            p.font.size = Pt(15)
            p.font.color.rgb = rgb("E2E8F0")
            p.space_after = Pt(8)


def add_flow_slide(prs, title: str, steps: list[tuple[str, str]]):
    from pptx.util import Inches, Pt

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb("0F172A")

    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(9), Inches(0.8))
    tp = title_box.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(28)
    tp.font.bold = True
    tp.font.color.rgb = rgb("F5F0E8")

    x_start = 0.5
    width = 1.55
    gap = 0.25
    for i, (step_title, step_desc) in enumerate(steps):
        x = x_start + i * (width + gap)
        box = slide.shapes.add_shape(1, Inches(x), Inches(2.0), Inches(width), Inches(2.2))
        box.fill.solid()
        box.fill.fore_color.rgb = rgb("1E293B")
        box.line.color.rgb = rgb("C9B896")

        num = slide.shapes.add_textbox(Inches(x + 0.1), Inches(2.15), Inches(0.4), Inches(0.4))
        np = num.text_frame.paragraphs[0]
        np.text = str(i + 1)
        np.font.size = Pt(22)
        np.font.bold = True
        np.font.color.rgb = rgb("C9B896")

        tb = slide.shapes.add_textbox(Inches(x + 0.1), Inches(2.55), Inches(width - 0.2), Inches(1.5))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = step_title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = rgb("F5F0E8")
        p2 = tf.add_paragraph()
        p2.text = step_desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = rgb("94A3B8")
        p2.space_before = Pt(6)

        if i < len(steps) - 1:
            arrow = slide.shapes.add_textbox(
                Inches(x + width + 0.02), Inches(2.85), Inches(0.2), Inches(0.3)
            )
            arrow.text_frame.paragraphs[0].text = "→"
            arrow.text_frame.paragraphs[0].font.size = Pt(18)
            arrow.text_frame.paragraphs[0].font.color.rgb = rgb("64748B")


def add_closing_slide(prs, title: str, bullets: list[str], image_path: Path | None = None):
    from pptx.util import Inches, Pt

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb("1E293B")

    if image_path and image_path.exists():
        slide.shapes.add_picture(str(image_path), Inches(6.0), Inches(1.0), width=Inches(3.5))

    box = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(5.2), Inches(4.5))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = rgb("F5F0E8")
    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(18)
        p.font.color.rgb = rgb("E2E8F0")
        p.space_before = Pt(12)


def build() -> Path:
    ensure_deps()
    copy_assets()

    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    hero = ASSETS / "ppt-hero-abu-dhabi.png"
    sim = ASSETS / "ppt-tenant-simulator.png"
    dash = ASSETS / "ppt-investment-dashboard.png"

    add_title_slide(
        prs,
        "The Sims 5 Real Estate",
        "AI intelligence for land, investment, and community decisions in Abu Dhabi",
        hero,
    )

    add_content_slide(
        prs,
        "The Problem",
        [
            "Investors see attractive assets — especially beach-facing buildings — but miss hidden risks.",
            "Low rental yield, high maintenance near the sea, and weak infrastructure can erase margins.",
            "Landlords lack a way to test how tenant life and management choices affect ROI over time.",
            "Decisions are often made on gut feel, not connected data and simulation.",
        ],
    )

    add_section_slide(
        prs,
        "Our Solution",
        "Two connected layers: score the asset, then simulate how you run it",
    )

    add_two_column_slide(
        prs,
        "One Product, Two Brains",
        "Investment Intelligence",
        [
            "Scores Abu Dhabi parcels and districts",
            "Predicts value, yield, and development potential",
            "Returns BUY · CONSIDER · DO NOT BUY",
            "Flags coastal and margin risks with warnings",
            "Live demo on the web dashboard today",
        ],
        "Tenant Simulator",
        [
            "Mini Sims-style building with AI residents",
            "Tenants talk, complain, pay rent, move in/out",
            "Landlord responds to repairs and complaints",
            "Updates reputation, occupancy, maintenance, ROI",
            "Partially built — visual panel coming next",
        ],
    )

    add_content_slide(
        prs,
        "Investment Dashboard",
        [
            "Uses synthetic Abu Dhabi starter-kit data (parcels, districts, transactions).",
            "Three local ML models estimate price/sqm, parcel value, and development potential.",
            "Transparent ROI rules turn predictions into a success score and recommendation.",
            "Investors get numbers plus plain-language reasons — not a black box.",
            "Disclaimer: triage aid for demo; not a real purchase decision.",
        ],
        dash,
    )

    add_content_slide(
        prs,
        "Beach-Facing Investor Risks We Surface",
        [
            "View premium → higher entry price and compressed gross yield.",
            "Salt air and humidity → higher maintenance and capex over time.",
            "Seasonal occupancy → churn and reputation volatility.",
            "Regulatory limits → harder redevelopment on waterfront parcels.",
            "Our coastal risk layer adjusts scores and adds due-diligence warnings.",
        ],
        hero,
    )

    add_content_slide(
        prs,
        "Tenant Simulator — How It Works",
        [
            "Demo building: 2 floors, 6 units, ~8 AI characters with unique personas.",
            "Residents propose actions; a game master validates what actually happens.",
            "Landlord handles repairs, complaints, rent issues, and move-ins.",
            "Simulation tracks satisfaction, trust, maintenance pressure, and budget.",
            "Connects operational reality to investor ROI — especially for premium coastal assets.",
        ],
        sim,
    )

    add_flow_slide(
        prs,
        "End-to-End Flow",
        [
            ("Data", "Synthetic CSVs\nDistricts · Parcels · Transactions"),
            ("ML Models", "Value · Yield ·\nDevelopment potential"),
            ("ROI Rules", "Margin · Score ·\nWarnings"),
            ("Simulator", "Tenant life ·\nLandlord choices"),
            ("Decision", "Invest & manage\nwith evidence"),
        ],
    )

    add_content_slide(
        prs,
        "Tech Stack",
        [
            "Next.js 14 + TypeScript + Tailwind — demo dashboard and API routes.",
            "Python ML pipeline — CatBoost & XGBoost models, ONNX export option.",
            "Agent harness — character, landlord, and game master agents with LLM hooks.",
            "Phaser (planned) — visual building panel for the tenant simulator.",
            "Local-first demo — runs without paid APIs; keys optional for AI memos.",
        ],
    )

    add_content_slide(
        prs,
        "3-Minute Demo Script",
        [
            "1. Frame the problem: beach views ≠ guaranteed returns.",
            "2. Open dashboard → Fetch recommendation for a waterfront parcel.",
            "3. Show BUY/CONSIDER score, margin, yield, and coastal warnings.",
            "4. Explain tenant sim: how bad management erodes ROI after purchase.",
            "5. Close: synthetic data, transparent logic, decision support not autopilot.",
        ],
        dash,
    )

    add_closing_slide(
        prs,
        "Thank You",
        [
            "The Sims 5 Real Estate",
            "Building the intelligence layer for land, investment, and communities.",
            "Questions?",
        ],
        sim,
    )

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    out_path = OUT_DIR / "The-Sims-5-Real-Estate-Presentation.pptx"
    prs.save(str(out_path))
    return out_path


if __name__ == "__main__":
    path = build()
    print(f"Saved: {path}")
